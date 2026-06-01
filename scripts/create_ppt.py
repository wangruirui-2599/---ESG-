#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
ESG Insight Valuator — 精美PPT生成器
=====================================
生成8-10页可用于展示汇报的专业PPT，包含图表嵌入和数据分析。

使用方式:
  python scripts/create_ppt.py
"""

import sys
from pathlib import Path
from datetime import datetime

import pandas as pd
import numpy as np

sys.path.insert(0, str(Path(__file__).parent.parent))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================================
# 颜色主题 — 深蓝金配色
# ============================================================================

NAVY = RGBColor(0x0B, 0x1D, 0x3A)       # 深海军蓝
DARK_BLUE = RGBColor(0x14, 0x2D, 0x5E)   # 暗蓝
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xAB)  # 强调蓝
GOLD = RGBColor(0xD4, 0xA8, 0x3C)         # 金色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xEC, 0xF0)
DARK_GRAY = RGBColor(0x50, 0x55, 0x60)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)


# ============================================================================
# 辅助函数
# ============================================================================

def add_slide_bg(slide, color=NAVY):
    """设置幻灯片纯色背景。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, opacity=None):
    """添加矩形色块。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """添加文本框。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_name="Microsoft YaHei"):
    """添加多行文本，每行可指定不同样式。lines: [(text, size, color, bold, align), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        text, size, color, bold, align = line
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
    return txBox


def add_gold_line(slide, left, top, width):
    """添加金色装饰线。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


def add_image_safe(slide, img_path, left, top, width, height=None):
    """安全添加图片，不存在则跳过。"""
    path = Path(img_path)
    if path.exists():
        if height:
            return slide.shapes.add_picture(str(path), left, top, width, height)
        else:
            return slide.shapes.add_picture(str(path), left, top, width)
    return None


# ============================================================================
# 幻灯片生成
# ============================================================================

def load_data():
    """加载最新数据。"""
    df = pd.read_parquet("data/processed/07_advice.parquet")
    latest = df.sort_values("report_year").groupby("stock_code").last().reset_index()
    return df, latest


def slide_01_title(prs):
    """第1页：封面。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_slide_bg(slide, NAVY)

    # 装饰色块
    add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(3.5), DARK_BLUE)
    add_gold_line(slide, Inches(1), Inches(3.6), Inches(3))

    # 主标题
    add_text_box(slide, Inches(1), Inches(1.2), Inches(8), Inches(1.5),
                 "ESG Insight Valuator", font_size=44, color=WHITE, bold=True)

    # 副标题
    add_text_box(slide, Inches(1), Inches(2.5), Inches(8), Inches(0.8),
                 "ESG智能估值分析系统 — 行业综合对比报告", font_size=22, color=GOLD)

    # 金线
    add_gold_line(slide, Inches(1), Inches(3.9), Inches(2))

    # 信息栏
    add_multi_text(slide, Inches(1), Inches(4.5), Inches(8), Inches(2.5), [
        (f"报告日期: {datetime.now().strftime('%Y年%m月%d日')}", 16, LIGHT_GRAY, False, PP_ALIGN.LEFT),
        ("", 10, WHITE, False, PP_ALIGN.LEFT),
        ("覆盖范围: 10大行业 · 10家A股上市公司", 16, LIGHT_GRAY, False, PP_ALIGN.LEFT),
        ("数据周期: 2019–2024 · 6年连续追踪", 16, LIGHT_GRAY, False, PP_ALIGN.LEFT),
        ("分析维度: ESG评分 · DCF估值 · 异常检测 · 风险传导", 16, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ])

    # 底部装饰
    add_shape_bg(slide, Inches(0), Inches(7.0), Inches(10), Inches(0.5), DARK_BLUE)
    add_text_box(slide, Inches(1), Inches(7.05), Inches(8), Inches(0.4),
                 "CONFIDENTIAL  |  ESG Insight Valuator v1.0", font_size=9, color=DARK_GRAY)


def slide_02_overview(prs):
    """第2页：项目概述。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, WHITE)

    # 顶部色条
    add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(0.08), GOLD)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "PROJECT OVERVIEW", font_size=14, color=GOLD, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.8),
                 "项目概述", font_size=28, color=NAVY, bold=True)
    add_gold_line(slide, Inches(0.8), Inches(1.5), Inches(1.5))

    # 核心能力卡片
    cards = [
        ("🌱", "ESG量化分析", "动态权重引擎\n行业传导分析\n趋势动量评估"),
        ("💰", "DCF多情景估值", "乐观/中性/悲观\n概率加权期望值\n情绪因子校准"),
        ("🔍", "财务异常检测", "LightGBM模型\n多维规则标签\n风险等级分类"),
        ("📊", "智能报告生成", "HTML/Markdown\n可视化图表\n投资建议输出"),
    ]

    for i, (icon, title, desc) in enumerate(cards):
        x = Inches(0.5 + i * 2.35)
        y = Inches(2.2)
        card = add_shape_bg(slide, x, y, Inches(2.1), Inches(3.2), LIGHT_GRAY)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(1.7), Inches(0.5),
                     icon, font_size=28, color=ACCENT_BLUE)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(1.7), Inches(0.4),
                     title, font_size=16, color=NAVY, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.3), Inches(1.7), Inches(1.7),
                     desc, font_size=11, color=DARK_GRAY)

    # 底部流程
    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(8.5), Inches(0.4),
                 "分析管线", font_size=14, color=NAVY, bold=True)
    steps = "数据加载 → 特征工程 → ESG量化 → 异常检测 → DCF估值 → 因子融合 → 投资建议 → 回测 → 报告"
    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(8.5), Inches(0.6),
                 steps, font_size=11, color=DARK_GRAY)
    add_shape_bg(slide, Inches(0.8), Inches(6.85), Inches(8.5), Inches(0.06), ACCENT_BLUE)

    add_text_box(slide, Inches(0.8), Inches(7.0), Inches(8.5), Inches(0.3),
                 "ESG Insight Valuator", font_size=9, color=DARK_GRAY)


def slide_03_data(prs, latest, df):
    """第3页：数据概览。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, NAVY)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), GOLD)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
                 "DATA COVERAGE", font_size=14, color=GOLD, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.7),
                 "数据覆盖与行业分布", font_size=28, color=WHITE, bold=True)
    add_gold_line(slide, Inches(0.8), Inches(1.4), Inches(1.5))

    # 指标卡片
    metrics = [
        ("10", "覆盖行业", ACCENT_BLUE),
        ("10", "标的公司", GREEN),
        ("6年", "追踪周期", GOLD),
        ("60条", "数据记录", ORANGE),
    ]
    for i, (val, label, clr) in enumerate(metrics):
        x = Inches(0.8 + i * 2.3)
        y = Inches(1.9)
        add_shape_bg(slide, x, y, Inches(2.0), Inches(1.3), DARK_BLUE)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(1.6), Inches(0.5),
                     val, font_size=32, color=clr, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(1.6), Inches(0.4),
                     label, font_size=13, color=LIGHT_GRAY)

    # 行业列表
    industries = sorted(latest["industry"].unique())
    n_cols = 5
    for i, ind in enumerate(industries):
        col = i % n_cols
        row = i // n_cols
        x = Inches(0.8 + col * 1.85)
        y = Inches(3.6 + row * 0.45)
        add_text_box(slide, x, y, Inches(1.7), Inches(0.35),
                     f"●  {ind}", font_size=12, color=LIGHT_GRAY)

    # ESG 改善亮点
    esg_2019 = df[df["report_year"] == 2019]["ESG_total"].mean()
    esg_2024 = df[df["report_year"] == 2024]["ESG_total"].mean()
    delta = esg_2024 - esg_2019

    add_shape_bg(slide, Inches(0.8), Inches(4.8), Inches(8.5), Inches(2.0), DARK_BLUE)
    add_multi_text(slide, Inches(1.2), Inches(5.0), Inches(7.5), Inches(1.6), [
        (f"ESG总分趋势: {esg_2019:.1f} (2019) → {esg_2024:.1f} (2024)  提升 +{delta:.1f} 分 (+{delta/esg_2019*100:.0f}%)", 16, GREEN, True, PP_ALIGN.LEFT),
        ("", 8, WHITE, False, PP_ALIGN.LEFT),
        ("· 全行业ESG总分6年持续提升，E维度改善最为显著（+9.4分），反映碳中和政策推动效应", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
        ("· 治理(G)维度普遍较高，社会(S)维度增长稳定，环境(E)维度行业分化明显", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
        ("· 10家公司中10家ESG趋势评级为'改善'或'显著改善'", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ])

    add_text_box(slide, Inches(0.8), Inches(7.0), Inches(8), Inches(0.3),
                 "ESG Insight Valuator", font_size=9, color=DARK_GRAY)


def slide_04_esg_ranking(prs, latest):
    """第4页：ESG评分排名。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), GOLD)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(4), Inches(0.5),
                 "ESG SCORE RANKING", font_size=13, color=GOLD, bold=True)
    add_text_box(slide, Inches(0.5), Inches(0.65), Inches(5), Inches(0.7),
                 "行业ESG评分排名", font_size=26, color=NAVY, bold=True)
    add_gold_line(slide, Inches(0.5), Inches(1.25), Inches(1.2))

    # 左侧图表
    add_image_safe(slide, "output/figures/industry_esg_ranking.png",
                   Inches(0.3), Inches(1.6), Inches(6.0), Inches(4.2))

    # 右侧分析
    top3 = latest.nlargest(3, "ESG_total")
    bot3 = latest.nsmallest(3, "ESG_total")

    add_shape_bg(slide, Inches(6.5), Inches(1.6), Inches(3.2), Inches(5.0), LIGHT_GRAY)
    add_text_box(slide, Inches(6.7), Inches(1.8), Inches(2.8), Inches(0.4),
                 "🏆 TOP 3", font_size=16, color=NAVY, bold=True)

    for i, (_, row) in enumerate(top3.iterrows()):
        add_text_box(slide, Inches(6.7), Inches(2.3 + i * 0.55), Inches(2.8), Inches(0.5),
                     f"{'🥇🥈🥉'[i]} {row['industry']}: {row['ESG_total']:.0f}分",
                     font_size=12, color=DARK_GRAY, bold=True)

    add_text_box(slide, Inches(6.7), Inches(4.1), Inches(2.8), Inches(0.4),
                 "⚠ 需关注", font_size=16, color=RED, bold=True)
    for i, (_, row) in enumerate(bot3.iterrows()):
        add_text_box(slide, Inches(6.7), Inches(4.6 + i * 0.55), Inches(2.8), Inches(0.5),
                     f"{row['industry']}: {row['ESG_total']:.0f}分",
                     font_size=12, color=DARK_GRAY)

    # 底部洞察
    add_shape_bg(slide, Inches(0.5), Inches(6.1), Inches(9.2), Inches(1.0), NAVY)
    add_multi_text(slide, Inches(0.8), Inches(6.2), Inches(8.5), Inches(0.8), [
        ("💡 关键发现", 14, GOLD, True, PP_ALIGN.LEFT),
        ("信息技术(77.6分)领跑全行业，采矿(50.7分)垫底。行业ESG标准差8.0分，分化显著。治理(G)维度行业间差异最小，环境(E)维度差异最大。", 11, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ])


def slide_05_radar(prs):
    """第5页：ESG雷达图对比。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, NAVY)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), GOLD)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.5),
                 "ESG RADAR COMPARISON", font_size=13, color=GOLD, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.65), Inches(8), Inches(0.7),
                 "行业ESG三维度雷达对比", font_size=26, color=WHITE, bold=True)
    add_gold_line(slide, Inches(0.8), Inches(1.2), Inches(1.5))

    # 多行业雷达图
    add_image_safe(slide, "output/figures/industry_multi_radar.png",
                   Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5))

    # 右侧示例
    add_image_safe(slide, "output/figures/esg_radar_600010.png",
                   Inches(6.3), Inches(1.5), Inches(3.2), Inches(3.2))

    add_image_safe(slide, "output/figures/esg_radar_600008.png",
                   Inches(6.3), Inches(3.9), Inches(3.2), Inches(3.2))

    add_shape_bg(slide, Inches(0.8), Inches(7.1), Inches(8.5), Inches(0.25), DARK_BLUE)
    add_text_box(slide, Inches(0.8), Inches(7.05), Inches(8.5), Inches(0.3),
                 "ESG Insight Valuator", font_size=9, color=DARK_GRAY)


def slide_06_trends(prs, df):
    """第6页：ESG时间趋势。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), GOLD)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(4), Inches(0.5),
                 "ESG TREND ANALYSIS", font_size=13, color=GOLD, bold=True)
    add_text_box(slide, Inches(0.5), Inches(0.65), Inches(5), Inches(0.7),
                 "ESG趋势演变 (2019–2024)", font_size=26, color=NAVY, bold=True)
    add_gold_line(slide, Inches(0.5), Inches(1.25), Inches(1.2))

    # 图表
    add_image_safe(slide, "output/figures/industry_esg_timeseries.png",
                   Inches(0.3), Inches(1.5), Inches(9.4), Inches(4.5))

    # 年度数据表
    years = sorted(df["report_year"].unique())
    y_data = []
    for y in years:
        yd = df[df["report_year"] == y]
        y_data.append((int(y), yd["E_score"].mean(), yd["S_score"].mean(),
                        yd["G_score"].mean(), yd["ESG_total"].mean()))

    add_text_box(slide, Inches(0.5), Inches(6.2), Inches(9), Inches(0.3),
                 "年度ESG均值变化", font_size=12, color=NAVY, bold=True)

    # 迷你表格
    table_data = [["Year", "E Score", "S Score", "G Score", "ESG Total"]] + \
                 [[f"{int(d[0])}", f"{d[1]:.1f}", f"{d[2]:.1f}", f"{d[3]:.1f}", f"{d[4]:.1f}"] for d in y_data]

    rows, cols = len(table_data), len(table_data[0])
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(6.45), Inches(9), Inches(0.5)).table

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = table_data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.name = "Microsoft YaHei"
                p.alignment = PP_ALIGN.CENTER
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = NAVY
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY


def slide_07_risk(prs):
    """第7页：风险传导分析。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, NAVY)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), GOLD)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.5),
                 "RISK CONTAGION ANALYSIS", font_size=13, color=GOLD, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.65), Inches(8), Inches(0.7),
                 "行业风险传导与投资吸引力", font_size=26, color=WHITE, bold=True)
    add_gold_line(slide, Inches(0.8), Inches(1.2), Inches(1.5))

    # 左侧：传导图
    add_image_safe(slide, "output/figures/industry_contagion.png",
                   Inches(0.2), Inches(1.4), Inches(5.8), Inches(3.0))

    # 右侧：吸引力矩阵
    add_image_safe(slide, "output/figures/industry_attractiveness.png",
                   Inches(6.2), Inches(1.4), Inches(3.5), Inches(3.0))

    # 分析文本
    add_shape_bg(slide, Inches(0.5), Inches(4.6), Inches(9.2), Inches(2.3), DARK_BLUE)
    add_multi_text(slide, Inches(0.8), Inches(4.8), Inches(8.5), Inches(2.0), [
        ("💡 风险与投资洞察", 16, GOLD, True, PP_ALIGN.LEFT),
        ("", 6, WHITE, False, PP_ALIGN.LEFT),
        ("🔴 高传导风险行业: 汽车(82分)、家电(59分) — 长供应链行业更易受上游波动冲击", 12, LIGHT_GRAY, False, PP_ALIGN.LEFT),
        ("🟢 低传导风险行业: 电力(0分)、采矿(0分)、食品饮料(0分) — 供应链独立性强", 12, LIGHT_GRAY, False, PP_ALIGN.LEFT),
        ("📊 投资吸引力: 信息技术、医药生物处于高ESG象限；采矿行业存在ESG改善驱动的估值修复机会", 12, LIGHT_GRAY, False, PP_ALIGN.LEFT),
        ("⚠ 银行行业G维度(91分)表现突出，但E维度(68分)相对较弱，需关注绿色信贷转型", 12, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ])


def slide_08_momentum(prs):
    """第8页：ESG质量与动量。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), GOLD)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(4), Inches(0.5),
                 "ESG QUALITY & MOMENTUM", font_size=13, color=GOLD, bold=True)
    add_text_box(slide, Inches(0.5), Inches(0.65), Inches(5), Inches(0.7),
                 "ESG质量与动量分析", font_size=26, color=NAVY, bold=True)
    add_gold_line(slide, Inches(0.5), Inches(1.25), Inches(1.2))

    # 气泡图
    add_image_safe(slide, "output/figures/industry_momentum_bubble.png",
                   Inches(0.3), Inches(1.5), Inches(9.4), Inches(4.3))

    # 底部洞察
    add_shape_bg(slide, Inches(0.5), Inches(6.0), Inches(9.2), Inches(1.0), NAVY)
    add_multi_text(slide, Inches(0.8), Inches(6.1), Inches(8.5), Inches(0.8), [
        ("💡 质量与动量洞察", 14, GOLD, True, PP_ALIGN.LEFT),
        ("右上象限（高质量+正动量）的行业具备最佳ESG投资属性。气泡大小反映财务健康度——气泡越大、异常风险越低。采矿行业动量最强（+0.049），显示ESG改善加速。", 11, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ])


def slide_09_insights(prs, latest, df):
    """第9页：核心发现与建议。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, NAVY)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), GOLD)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.5),
                 "KEY FINDINGS & RECOMMENDATIONS", font_size=13, color=GOLD, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.65), Inches(8), Inches(0.7),
                 "核心发现与投资建议", font_size=26, color=WHITE, bold=True)
    add_gold_line(slide, Inches(0.8), Inches(1.2), Inches(1.5))

    # 五条核心发现
    findings = [
        ("📈", "ESG整体向好", "全行业ESG均值从2019年56.4分提升至2024年66.6分(+18%)，碳中和政策推动效果显著"),
        ("🏆", "行业分化明显", "信息技术(77.6分) vs 采矿(50.7分)，差距达26.9分。环境(E)维度是主要分化来源"),
        ("📊", "治理维度最均衡", "G维度行业间标准差最小，中国企业治理水平整体提升，董事会独立性增强"),
        ("⚠", "供应链风险传导", "汽车行业受4个上游行业影响，传导评分82分，需关注供应链ESG风险管理"),
        ("💡", "ESG改善驱动估值", "采矿行业动量最高(+0.049)，存在ESG改善驱动的估值修复投资机会"),
    ]

    for i, (icon, title, desc) in enumerate(findings):
        y = Inches(1.6 + i * 1.1)
        add_shape_bg(slide, Inches(0.8), y, Inches(8.5), Inches(0.95), DARK_BLUE)
        add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(0.4), Inches(0.4),
                     icon, font_size=20, color=GOLD)
        add_text_box(slide, Inches(1.5), y + Inches(0.05), Inches(2.0), Inches(0.35),
                     title, font_size=15, color=WHITE, bold=True)
        add_text_box(slide, Inches(1.5), y + Inches(0.45), Inches(7.5), Inches(0.45),
                     desc, font_size=10.5, color=LIGHT_GRAY)

    add_text_box(slide, Inches(0.8), Inches(7.1), Inches(8), Inches(0.3),
                 "ESG Insight Valuator", font_size=9, color=DARK_GRAY)


def slide_10_thanks(prs):
    """第10页：致谢。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, NAVY)

    add_shape_bg(slide, Inches(0), Inches(0), Inches(10), Inches(3.2), DARK_BLUE)
    add_gold_line(slide, Inches(3), Inches(3.4), Inches(4))

    add_text_box(slide, Inches(1), Inches(1.5), Inches(8), Inches(1),
                 "THANK YOU", font_size=48, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.3), Inches(8), Inches(0.6),
                 "感谢聆听 · 欢迎交流", font_size=20, color=GOLD,
                 alignment=PP_ALIGN.CENTER)

    add_multi_text(slide, Inches(2), Inches(4.2), Inches(6), Inches(2.5), [
        ("ESG Insight Valuator v1.0", 18, WHITE, True, PP_ALIGN.CENTER),
        ("", 10, WHITE, False, PP_ALIGN.CENTER),
        ("ESG智能估值分析系统", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER),
        ("", 10, WHITE, False, PP_ALIGN.CENTER),
        ("数据来源: 10家A股上市公司 2019-2024年度ESG与财务数据", 11, DARK_GRAY, False, PP_ALIGN.CENTER),
        ("分析方法: DCF多情景估值 · LightGBM异常检测 · 行业动态权重 · 风险传导网络", 11, DARK_GRAY, False, PP_ALIGN.CENTER),
        ("", 10, WHITE, False, PP_ALIGN.CENTER),
        ("⚠ 本报告仅供研究参考，不构成投资建议", 10, RED, False, PP_ALIGN.CENTER),
    ])

    add_shape_bg(slide, Inches(0), Inches(7.0), Inches(10), Inches(0.5), DARK_BLUE)
    add_text_box(slide, Inches(1), Inches(7.05), Inches(8), Inches(0.4),
                 "CONFIDENTIAL  |  ESG Insight Valuator", font_size=9, color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER)


# ============================================================================
# 主入口
# ============================================================================

def main():
    print("Loading data...")
    df, latest = load_data()

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Creating slides...")
    slide_01_title(prs)
    slide_02_overview(prs)
    slide_03_data(prs, latest, df)
    slide_04_esg_ranking(prs, latest)
    slide_05_radar(prs)
    slide_06_trends(prs, df)
    slide_07_risk(prs)
    slide_08_momentum(prs)
    slide_09_insights(prs, latest, df)
    slide_10_thanks(prs)

    output_path = "output/reports/ESG_Insight_Valuator_Presentation.pptx"
    prs.save(output_path)
    print(f"PPT saved: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
